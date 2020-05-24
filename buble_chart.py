from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = XyChartData()

series_1 = chart_data.add_series('Model 1')
series_1.add_data_point(0.7, 2.7)
series_1.add_data_point(1.8, 3.2)
series_1.add_data_point(2.6, 0.8)

series_2 = chart_data.add_series('Model 2')
series_2.add_data_point(1.3, 3.7)
series_2.add_data_point(2.7, 2.3)
series_2.add_data_point(1.6, 1.8)

series_2 = chart_data.add_series('Model 3')
series_2.add_data_point(3.3, 1.7)
series_2.add_data_point(2.9, 6.3)
series_2.add_data_point(7.6, 3.8)

series_2 = chart_data.add_series('Model 4')
series_2.add_data_point(2.3, 6.7)
series_2.add_data_point(1.9, 3.3)
series_2.add_data_point(1.4, 1.8)



x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
).chart

prs.save('chart-02.pptx')