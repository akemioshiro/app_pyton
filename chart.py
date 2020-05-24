import chart as chart
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest', 'Center']
chart_data.add_series('Visão Empresa', [19.2, 21.4, 16.7, 10])
chart_data.add_series('Visão Area', [22.3, 28.6, 15.2, 12])



# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False

chart.has_title = True
chart.chart_title.has_text_frame = True
chart.chart_title.text_frame.text = 'Response'

prs.save('chart-01.pptx')